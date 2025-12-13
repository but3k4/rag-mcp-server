"""
Text extraction and section-aware chunking for supported file types.

Each parser returns a list of chunk dicts with three keys:
  - text: the chunk content (stripped)
  - section: section title (header, slide title, page label, or empty string)
  - chunk_index: sequential integer within the file
"""

from __future__ import annotations

import csv
import logging
import re
from typing import TYPE_CHECKING, Any

from rag.errors import RagError

if TYPE_CHECKING:
    from pathlib import Path

logger = logging.getLogger(__name__)
logger.addHandler(logging.NullHandler())

Chunk = dict[str, Any]

SUPPORTED_EXTENSIONS: frozenset[str] = frozenset(
    {".txt", ".md", ".pdf", ".docx", ".pptx", ".csv", ".xlsx", ".xml"}
)

DEFAULT_CHUNK_SIZE = 1400  # characters
DEFAULT_CHUNK_OVERLAP = 150


class ParseError(RagError):
    """Raised when a file cannot be parsed into chunks."""


def chunk_file(
    path: Path,
    chunk_size: int = DEFAULT_CHUNK_SIZE,
    chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
) -> list[Chunk]:
    """
    Parse a file and return section-aware text chunks.

    Args:
        path: Absolute path to the file to parse.
        chunk_size: Maximum characters per chunk.
        chunk_overlap: Characters of overlap between consecutive chunks.

    Returns:
        List of chunk dicts, each with text, section, and chunk_index keys.
        Returns an empty list for files with no extractable text.

    Raises:
        ParseError: If the file type is unsupported or extraction fails.
    """

    suffix = path.suffix.lower()
    dispatch = {
        ".txt": _chunk_text,
        ".md": _chunk_markdown,
        ".pdf": _chunk_pdf,
        ".docx": _chunk_docx,
        ".pptx": _chunk_pptx,
        ".csv": _chunk_csv,
        ".xlsx": _chunk_xlsx,
        ".xml": _chunk_xml,
    }

    fn = dispatch.get(suffix)
    if fn is None:
        raise ParseError(f"Unsupported file type: {suffix}")

    return fn(path, chunk_size, chunk_overlap)


def _make_chunk(text: str, section: str, index: int) -> Chunk:
    """Build a Chunk dict with stripped text, section title, and index."""

    return {"text": text.strip(), "section": section, "chunk_index": index}


def _split_text(text: str, chunk_size: int, chunk_overlap: int) -> list[str]:
    """
    Split text into overlapping chunks at paragraph boundaries.

    Prefers splitting on double-newlines rather than mid-sentence. Each chunk
    is at most chunk_size characters. Consecutive chunks overlap by
    chunk_overlap characters to preserve context at boundaries.

    Args:
        text: The text to split.
        chunk_size: Maximum characters per chunk.
        chunk_overlap: Characters to repeat at the start of each new chunk.

    Returns:
        List of chunk strings. Returns an empty list for blank input.
    """

    if not text.strip():
        return []

    paragraphs = re.split(r"\n\n+", text)
    chunks: list[str] = []
    current = ""

    for para in paragraphs:
        if len(current) + len(para) > chunk_size and current:
            chunks.append(current.strip())
            current = (
                current[-chunk_overlap:] + "\n\n" + para if chunk_overlap else para
            )
        else:
            current = (current + "\n\n" + para).strip() if current else para

    if current.strip():
        chunks.append(current.strip())
    return chunks or [text.strip()]


def _chunk_text(path: Path, chunk_size: int, chunk_overlap: int) -> list[Chunk]:
    """
    Read a plain text file and split it into overlapping unsectioned chunks.

    Args:
        path: Path to the .txt file.
        chunk_size: Maximum characters per chunk.
        chunk_overlap: Characters of overlap between consecutive chunks.

    Returns:
        List of chunk dicts with empty section titles.

    Raises:
        ParseError: If the file cannot be read.
    """

    try:
        text = path.read_text(encoding="utf-8", errors="replace")
    except OSError as exc:
        raise ParseError(f"Cannot read {path}: {exc}") from exc
    return [
        _make_chunk(c, "", i)
        for i, c in enumerate(_split_text(text, chunk_size, chunk_overlap))
    ]


def _chunk_markdown(path: Path, chunk_size: int, chunk_overlap: int) -> list[Chunk]:
    """
    Split a Markdown file on ATX headers (# through ####) into sections.

    Each header line becomes the section title for the text that follows it,
    up to the next header of depth 1–4. Sections longer than chunk_size are
    further split into overlapping sub-chunks.

    Args:
        path: Path to the .md file.
        chunk_size: Maximum characters per chunk.
        chunk_overlap: Characters of overlap between consecutive chunks.

    Returns:
        List of chunk dicts keyed by header section.

    Raises:
        ParseError: If the file cannot be read.
    """

    try:
        text = path.read_text(encoding="utf-8", errors="replace")
    except OSError as exc:
        raise ParseError(f"Cannot read {path}: {exc}") from exc

    sections = re.split(r"(?=^#{1,4} )", text, flags=re.MULTILINE)
    chunks: list[Chunk] = []

    for section in sections:
        if not section.strip():
            continue

        header_match = re.match(r"^(#{1,4} .+)", section)
        section_title = header_match.group(1).strip() if header_match else ""

        for sub in _split_text(section, chunk_size, chunk_overlap):
            chunks.append(_make_chunk(sub, section_title, len(chunks)))
    return chunks


def _chunk_pdf(path: Path, chunk_size: int, chunk_overlap: int) -> list[Chunk]:
    """
    Extract text from each PDF page as a separate section.

    Args:
        path: Path to the PDF file.
        chunk_size: Maximum characters per chunk.
        chunk_overlap: Characters of overlap between consecutive chunks.

    Returns:
        List of chunk dicts, one or more per non-blank page.

    Raises:
        ParseError: If pypdf is not installed or the file cannot be parsed.
    """

    try:
        # I hate to import things inside functions but pypdf is a large
        # dependency and this keeps it out of the way for users who don't need
        # PDF support.
        import pypdf  # noqa: PLC0415
        from pypdf.errors import PyPdfError  # noqa: PLC0415
    except ImportError as exc:
        raise ParseError("pypdf is not installed") from exc

    try:
        with path.open("rb") as fh:
            reader = pypdf.PdfReader(fh)
            chunks: list[Chunk] = []

            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if not text.strip():
                    continue

                section = f"Page {i + 1}"

                for sub in _split_text(text, chunk_size, chunk_overlap):
                    chunks.append(_make_chunk(sub, section, len(chunks)))
    except (PyPdfError, OSError) as exc:
        raise ParseError(f"Cannot parse PDF {path}: {exc}") from exc
    return chunks


def _chunk_docx(path: Path, chunk_size: int, chunk_overlap: int) -> list[Chunk]:
    """
    Extract text from a Word document, splitting on Heading styles.

    Args:
        path: Path to the .docx file.
        chunk_size: Maximum characters per chunk.
        chunk_overlap: Characters of overlap between consecutive chunks.

    Returns:
        List of chunk dicts keyed by heading section.

    Raises:
        ParseError: If python-docx is not installed or parsing fails.
    """

    try:
        # Again, importing inside the function to avoid the dependency for
        # users who don't need DOCX support.
        import docx  # noqa: PLC0415
        from docx.opc.exceptions import OpcError  # noqa: PLC0415
    except ImportError as exc:
        raise ParseError("python-docx is not installed") from exc

    try:
        doc = docx.Document(str(path))
        chunks: list[Chunk] = []
        current_section = ""
        buffer: list[str] = []

        for para in doc.paragraphs:
            if para.style is not None and para.style.name.startswith("Heading"):
                if buffer:
                    for sub in _split_text(
                        "\n\n".join(buffer), chunk_size, chunk_overlap
                    ):
                        chunks.append(_make_chunk(sub, current_section, len(chunks)))
                    buffer.clear()
                current_section = para.text
            elif para.text.strip():
                buffer.append(para.text)

        if buffer:
            for sub in _split_text("\n\n".join(buffer), chunk_size, chunk_overlap):
                chunks.append(_make_chunk(sub, current_section, len(chunks)))

    except (OpcError, KeyError, ValueError, OSError) as exc:
        raise ParseError(f"Cannot parse DOCX {path}: {exc}") from exc
    return chunks


def _chunk_pptx(path: Path, chunk_size: int, chunk_overlap: int) -> list[Chunk]:
    """
    Extract text from each slide as a separate section.

    Args:
        path: Path to the .pptx file.
        chunk_size: Unused. Each slide is stored as a single chunk.
        chunk_overlap: Unused. Each slide is stored as a single chunk.

    Returns:
        List of chunk dicts, one per non-blank slide.

    Raises:
        ParseError: If python-pptx is not installed or parsing fails.
    """

    try:
        # Again, importing inside the function to avoid the dependency for
        # users who don't need PPTX support.
        from pptx import Presentation  # noqa: PLC0415
        from pptx.exc import PythonPptxError  # noqa: PLC0415
    except ImportError as exc:
        raise ParseError("python-pptx is not installed") from exc

    try:
        prs = Presentation(str(path))
        chunks: list[Chunk] = []
        for i, slide in enumerate(prs.slides):
            title = ""
            texts: list[str] = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                if shape.is_placeholder and shape.placeholder_format.idx == 0:
                    title = shape.text_frame.text
                else:
                    texts.append(shape.text_frame.text)
            content = "\n".join(t for t in texts if t.strip())
            combined = f"{title}\n{content}".strip() if title else content

            if combined:
                section = f"Slide {i + 1}: {title}" if title else f"Slide {i + 1}"
                chunks.append(_make_chunk(combined, section, len(chunks)))
    except (PythonPptxError, KeyError, ValueError, OSError) as exc:
        raise ParseError(f"Cannot parse PPTX {path}: {exc}") from exc
    return chunks


def _chunk_csv(path: Path, chunk_size: int, chunk_overlap: int) -> list[Chunk]:
    """
    Flatten a CSV file into pipe-delimited rows and split into chunks.

    Each row is joined with "|" between cells and rows are separated by
    newlines before being passed to the overlap-chunker. No header handling.

    Args:
        path: Path to the .csv file.
        chunk_size: Maximum characters per chunk.
        chunk_overlap: Characters of overlap between consecutive chunks.

    Returns:
        List of chunk dicts with empty section titles.

    Raises:
        ParseError: If the file cannot be read or parsed.
    """

    try:
        with path.open(encoding="utf-8", errors="replace", newline="") as fh:
            reader = csv.reader(fh)
            rows = [" | ".join(row) for row in reader]
    except OSError as exc:
        raise ParseError(f"Cannot read CSV {path}: {exc}") from exc
    except csv.Error as exc:
        raise ParseError(f"Cannot parse CSV {path}: {exc}") from exc

    text = "\n".join(rows)
    return [
        _make_chunk(c, "", i)
        for i, c in enumerate(_split_text(text, chunk_size, chunk_overlap))
    ]


def _chunk_xlsx(path: Path, chunk_size: int, chunk_overlap: int) -> list[Chunk]:
    """
    Flatten all worksheets in an XLSX file into pipe-delimited rows.

    Rows from every sheet are concatenated without sheet separators and
    empty cells become empty strings. Formulas are read as their cached
    values (data_only=True).

    Args:
        path: Path to the .xlsx file.
        chunk_size: Maximum characters per chunk.
        chunk_overlap: Characters of overlap between consecutive chunks.

    Returns:
        List of chunk dicts with empty section titles.

    Raises:
        ParseError: If openpyxl is not installed or the file cannot be parsed.
    """

    try:
        # Again, importing inside the function to avoid the dependency for
        # users who don't need XLSX support.
        import openpyxl  # noqa: PLC0415
        from openpyxl.utils.exceptions import InvalidFileException  # noqa: PLC0415
    except ImportError as exc:
        raise ParseError("openpyxl is not installed") from exc

    import zipfile  # noqa: PLC0415

    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts: list[str] = []

        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                parts.append(" | ".join(cells))
        wb.close()
    except (InvalidFileException, zipfile.BadZipFile, OSError) as exc:
        raise ParseError(f"Cannot parse XLSX {path}: {exc}") from exc

    text = "\n".join(parts)
    return [
        _make_chunk(c, "", i)
        for i, c in enumerate(_split_text(text, chunk_size, chunk_overlap))
    ]


def _chunk_xml(path: Path, chunk_size: int, chunk_overlap: int) -> list[Chunk]:
    """
    Extract text from an XML file, using direct child tag names as sections.

    Each direct child of the root element is treated as a section. Text is
    extracted from the full subtree of each child and split into overlapping
    chunks. For flat documents with no children, all root text is returned as
    a single unsectioned chunk group.

    Args:
        path: Path to the XML file.
        chunk_size: Maximum characters per chunk.
        chunk_overlap: Characters of overlap between consecutive chunks.

    Returns:
        List of chunk dicts, one or more per non-blank section.

    Raises:
        ParseError: If lxml is not installed or the file cannot be parsed.
    """

    try:
        # Again, importing inside the function to avoid the dependency for
        # users who don't need XML support.
        from lxml import etree  # noqa: PLC0415
    except ImportError as exc:
        raise ParseError("lxml is not installed") from exc

    try:
        parser = etree.XMLParser(
            resolve_entities=False, no_network=True, load_dtd=False
        )
        tree = etree.parse(str(path), parser)
        root = tree.getroot()
        chunks: list[Chunk] = []

        for child in root:
            if not isinstance(child.tag, str):
                continue

            tag = etree.QName(child.tag).localname
            text = " ".join(child.itertext()).strip()

            if not text:
                continue

            for sub in _split_text(text, chunk_size, chunk_overlap):
                chunks.append(_make_chunk(sub, tag, len(chunks)))
        if not chunks:
            text = " ".join(root.itertext()).strip()
            chunks = [
                _make_chunk(sub, "", i)
                for i, sub in enumerate(_split_text(text, chunk_size, chunk_overlap))
            ]
    except (etree.LxmlError, OSError) as exc:
        raise ParseError(f"Cannot parse XML {path}: {exc}") from exc

    return chunks
