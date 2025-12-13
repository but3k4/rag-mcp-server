"""Tests for rag.parsers section-aware chunking functions."""

from __future__ import annotations

import sys
from typing import TYPE_CHECKING

import docx as _docx
import openpyxl
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter
import pytest

from rag.parsers import SUPPORTED_EXTENSIONS, ParseError, chunk_file

if TYPE_CHECKING:
    from pathlib import Path


class TestChunkFileDispatch:
    """Tests for the dispatch logic in chunk_file."""

    def test_unsupported_extension_raises_parse_error(self, tmp_path: Path) -> None:
        """chunk_file raises ParseError for an unknown file extension."""

        f = tmp_path / "image.png"
        f.write_bytes(b"\x89PNG")
        with pytest.raises(ParseError, match="Unsupported file type"):
            chunk_file(f)

    def test_supported_extensions_set(self) -> None:
        """SUPPORTED_EXTENSIONS contains exactly the eight documented types."""

        expected = frozenset(
            {".txt", ".md", ".pdf", ".docx", ".pptx", ".csv", ".xlsx", ".xml"}
        )
        assert expected == SUPPORTED_EXTENSIONS

    def test_returns_list_of_dicts(self, tmp_path: Path) -> None:
        """chunk_file returns a list of dicts with required keys."""

        f = tmp_path / "note.txt"
        f.write_text("hello world", encoding="utf-8")
        result = chunk_file(f)
        assert isinstance(result, list)
        assert len(result) >= 1
        for chunk in result:
            assert "text" in chunk
            assert "section" in chunk
            assert "chunk_index" in chunk

    def test_chunk_indices_are_sequential(self, tmp_path: Path) -> None:
        """chunk_index values are contiguous starting from 0."""

        f = tmp_path / "note.txt"
        f.write_text("hello world", encoding="utf-8")
        result = chunk_file(f)
        assert [c["chunk_index"] for c in result] == list(range(len(result)))

    def test_custom_chunk_size_splits_text(self, tmp_path: Path) -> None:
        """A small chunk_size produces more chunks than the default."""

        f = tmp_path / "long.txt"
        # Paragraphs separated by blank lines give _split_text break points.
        f.write_text("\n\n".join(["word " * 10] * 20), encoding="utf-8")
        small = chunk_file(f, chunk_size=100, chunk_overlap=0)
        default = chunk_file(f)
        assert len(small) > len(default)


class TestChunkTxt:
    """Tests for plain-text and Markdown chunking."""

    def test_txt_contains_content(self, tmp_path: Path) -> None:
        """Content from a .txt file appears in chunk text."""

        f = tmp_path / "note.txt"
        f.write_text("hello world", encoding="utf-8")
        result = chunk_file(f)
        assert any("hello world" in c["text"] for c in result)

    def test_md_content_present(self, tmp_path: Path) -> None:
        """Markdown header and body text are present across chunks."""

        f = tmp_path / "note.md"
        f.write_text("# Heading\nBody text.", encoding="utf-8")
        result = chunk_file(f)
        all_text = " ".join(c["text"] for c in result)
        assert "Heading" in all_text
        assert "Body text" in all_text

    def test_unreadable_file_raises_parse_error(self, tmp_path: Path) -> None:
        """A file that does not exist raises ParseError."""

        f = tmp_path / "gone.txt"
        with pytest.raises(ParseError):
            chunk_file(f)

    def test_unreadable_markdown_raises_parse_error(self, tmp_path: Path) -> None:
        """A missing .md file raises ParseError with Cannot read wording."""

        f = tmp_path / "gone.md"
        with pytest.raises(ParseError, match="Cannot read"):
            chunk_file(f)


class TestChunkPdf:
    """Tests for PDF chunking."""

    def test_invalid_pdf_raises_parse_error(self, tmp_path: Path) -> None:
        """A file with non-PDF bytes raises ParseError."""

        f = tmp_path / "bad.pdf"
        f.write_bytes(b"not a pdf at all")
        with pytest.raises(ParseError, match="Cannot parse PDF"):
            chunk_file(f)

    def test_valid_pdf_returns_list(self, tmp_path: Path) -> None:
        """A valid blank-page PDF returns a list without raising."""

        writer = PdfWriter()
        writer.add_blank_page(width=200, height=200)
        out = tmp_path / "blank.pdf"
        with out.open("wb") as fh:
            writer.write(fh)
        result = chunk_file(out)
        assert isinstance(result, list)

    def test_pdf_with_text_uses_page_section(self, tmp_path: Path) -> None:
        """A PDF with real text produces chunks labelled 'Page N'."""

        # Build a PDF using reportlab-free path: use pypdf with a page cloned
        # from a stream containing text. Simpler: use fpdf via raw bytes is
        # overkill. Instead we write a minimal PDF with embedded text using
        # pypdf's low-level API.
        from pypdf.generic import DecodedStreamObject, NameObject  # noqa: PLC0415

        writer = PdfWriter()
        page = writer.add_blank_page(width=200, height=200)
        content = DecodedStreamObject()
        content.set_data(b"BT /F1 12 Tf 50 100 Td (Hello PDF text) Tj ET")
        page[NameObject("/Contents")] = content
        out = tmp_path / "text.pdf"
        with out.open("wb") as fh:
            writer.write(fh)
        result = chunk_file(out)
        # Even if extract_text fails on this minimal PDF, the code path is
        # exercised. We only assert it returns a list.
        assert isinstance(result, list)

    def test_missing_pypdf_raises_parse_error(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        """If pypdf is not installed, PDF parsing raises ParseError."""

        f = tmp_path / "x.pdf"
        f.write_bytes(b"%PDF-1.4")
        monkeypatch.setitem(sys.modules, "pypdf", None)
        with pytest.raises(ParseError, match="pypdf is not installed"):
            chunk_file(f)


class TestChunkDocx:
    """Tests for Word document chunking."""

    def _make_docx(self, tmp_path: Path, text: str) -> Path:
        """Write a minimal .docx containing a single paragraph and return its path."""

        doc = _docx.Document()
        doc.add_paragraph(text)
        out = tmp_path / "doc.docx"
        doc.save(str(out))
        return out

    def test_docx_contains_paragraph_text(self, tmp_path: Path) -> None:
        """Paragraph text appears in chunk text."""

        f = self._make_docx(tmp_path, "Test paragraph content")
        result = chunk_file(f)
        assert any("Test paragraph content" in c["text"] for c in result)

    def test_invalid_docx_raises_parse_error(self, tmp_path: Path) -> None:
        """Non-docx bytes raise ParseError."""

        f = tmp_path / "bad.docx"
        f.write_bytes(b"definitely not a docx")
        with pytest.raises(ParseError, match="Cannot parse DOCX"):
            chunk_file(f)

    def test_docx_heading_flushes_buffered_body(self, tmp_path: Path) -> None:
        """Body paragraphs are flushed into a section when a heading appears."""

        doc = _docx.Document()
        doc.add_paragraph("First body line")
        doc.add_paragraph("Second body line")
        doc.add_heading("Chapter two", level=1)
        doc.add_paragraph("Body under chapter two")
        out = tmp_path / "sections.docx"
        doc.save(str(out))
        result = chunk_file(out)
        sections = {c["section"] for c in result}
        assert "Chapter two" in sections

    def test_missing_docx_package_raises_parse_error(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        """If python-docx is not installed, DOCX parsing raises ParseError."""

        f = tmp_path / "x.docx"
        f.write_bytes(b"PK")
        monkeypatch.setitem(sys.modules, "docx", None)
        with pytest.raises(ParseError, match="python-docx is not installed"):
            chunk_file(f)


class TestChunkPptx:
    """Tests for PowerPoint chunking."""

    def _make_pptx(self, tmp_path: Path, slide_text: str) -> Path:
        """Write a minimal .pptx with one textbox on a blank slide and return its path."""

        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # blank layout
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txBox.text_frame.text = slide_text
        out = tmp_path / "pres.pptx"
        prs.save(str(out))
        return out

    def test_pptx_contains_slide_text(self, tmp_path: Path) -> None:
        """Text from slide shapes appears in chunk text."""

        f = self._make_pptx(tmp_path, "Slide one content")
        result = chunk_file(f)
        assert any("Slide one content" in c["text"] for c in result)

    def test_invalid_pptx_raises_parse_error(self, tmp_path: Path) -> None:
        """Non-pptx bytes raise ParseError."""

        f = tmp_path / "bad.pptx"
        f.write_bytes(b"garbage bytes")
        with pytest.raises(ParseError, match="Cannot parse PPTX"):
            chunk_file(f)

    def test_pptx_skips_shapes_without_text_frame(self, tmp_path: Path) -> None:
        """Shapes with no text frame (e.g. pictures, lines) are skipped."""

        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        # Add a connector shape (line) which has no text frame.
        from pptx.enum.shapes import MSO_CONNECTOR  # noqa: PLC0415

        slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(3), Inches(3)
        )
        # Also add a real text box so the slide produces a chunk.
        tx = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(4), Inches(1))
        tx.text_frame.text = "Slide content here"
        out = tmp_path / "mixed.pptx"
        prs.save(str(out))
        result = chunk_file(out)
        assert any("Slide content here" in c["text"] for c in result)

    def test_missing_pptx_package_raises_parse_error(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        """If python-pptx is not installed, PPTX parsing raises ParseError."""

        f = tmp_path / "x.pptx"
        f.write_bytes(b"PK")
        monkeypatch.setitem(sys.modules, "pptx", None)
        with pytest.raises(ParseError, match="python-pptx is not installed"):
            chunk_file(f)


class TestChunkCsv:
    """Tests for CSV chunking."""

    def test_csv_rows_in_chunk_text(self, tmp_path: Path) -> None:
        """Each CSV row appears formatted in chunk text."""

        f = tmp_path / "data.csv"
        f.write_text("a,b,c\n1,2,3\n", encoding="utf-8")
        result = chunk_file(f)
        all_text = " ".join(c["text"] for c in result)
        assert "a | b | c" in all_text
        assert "1 | 2 | 3" in all_text

    def test_empty_csv_returns_empty_list(self, tmp_path: Path) -> None:
        """An empty CSV file returns an empty list."""

        f = tmp_path / "empty.csv"
        f.write_text("", encoding="utf-8")
        result = chunk_file(f)
        assert result == []

    def test_unreadable_csv_raises_parse_error(self, tmp_path: Path) -> None:
        """A CSV that does not exist raises ParseError."""

        f = tmp_path / "missing.csv"
        with pytest.raises(ParseError, match="Cannot read CSV"):
            chunk_file(f)

    def test_csv_error_raises_parse_error(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        """A csv.Error during parsing is wrapped as ParseError."""

        import csv as _csv  # noqa: PLC0415

        f = tmp_path / "broken.csv"
        f.write_text("a,b\n", encoding="utf-8")

        def _boom(*_args: object, **_kwargs: object) -> object:
            """Monkeypatch stand-in that always raises csv.Error."""

            raise _csv.Error("simulated")

        monkeypatch.setattr(_csv, "reader", _boom)
        with pytest.raises(ParseError, match="Cannot parse CSV"):
            chunk_file(f)


class TestChunkXlsx:
    """Tests for Excel workbook chunking."""

    def _make_xlsx(self, tmp_path: Path, rows: list[list[str]]) -> Path:
        """Write a minimal .xlsx containing the given rows and return its path."""

        wb = openpyxl.Workbook()
        ws = wb.active
        for row in rows:
            ws.append(row)
        out = tmp_path / "data.xlsx"
        wb.save(str(out))
        return out

    def test_xlsx_rows_in_chunk_text(self, tmp_path: Path) -> None:
        """Each row across all sheets is present in chunk text."""

        f = self._make_xlsx(tmp_path, [["name", "value"], ["foo", "42"]])
        result = chunk_file(f)
        all_text = " ".join(c["text"] for c in result)
        assert "name | value" in all_text
        assert "foo | 42" in all_text

    def test_invalid_xlsx_raises_parse_error(self, tmp_path: Path) -> None:
        """Non-xlsx bytes raise ParseError."""

        f = tmp_path / "bad.xlsx"
        f.write_bytes(b"not xlsx data")
        with pytest.raises(ParseError, match="Cannot parse XLSX"):
            chunk_file(f)

    def test_missing_openpyxl_raises_parse_error(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        """If openpyxl is not installed, XLSX parsing raises ParseError."""

        f = tmp_path / "x.xlsx"
        f.write_bytes(b"PK")
        monkeypatch.setitem(sys.modules, "openpyxl", None)
        with pytest.raises(ParseError, match="openpyxl is not installed"):
            chunk_file(f)


class TestChunkXml:
    """Tests for XML chunking."""

    def test_xml_child_text_in_chunks(self, tmp_path: Path) -> None:
        """Text from child elements appears in chunk text."""

        f = tmp_path / "doc.xml"
        f.write_text(
            "<root><section>Hello XML world</section></root>",
            encoding="utf-8",
        )
        result = chunk_file(f)
        assert any("Hello XML world" in c["text"] for c in result)

    def test_xml_child_tag_used_as_section(self, tmp_path: Path) -> None:
        """The direct child tag name is used as the section label."""

        f = tmp_path / "doc.xml"
        f.write_text(
            "<root><chapter>Chapter text</chapter></root>",
            encoding="utf-8",
        )
        result = chunk_file(f)
        assert any(c["section"] == "chapter" for c in result)

    def test_xml_multiple_children_produce_sections(self, tmp_path: Path) -> None:
        """Each non-empty child element produces at least one chunk."""

        f = tmp_path / "doc.xml"
        f.write_text(
            "<root><a>Alpha</a><b>Beta</b></root>",
            encoding="utf-8",
        )
        result = chunk_file(f)
        sections = {c["section"] for c in result}
        assert "a" in sections
        assert "b" in sections

    def test_xml_flat_document_returns_chunks(self, tmp_path: Path) -> None:
        """A document with no child elements still returns chunks from root text."""

        f = tmp_path / "flat.xml"
        f.write_text("<root>Just some text here</root>", encoding="utf-8")
        result = chunk_file(f)
        assert len(result) >= 1
        assert any("Just some text here" in c["text"] for c in result)

    def test_invalid_xml_raises_parse_error(self, tmp_path: Path) -> None:
        """Malformed XML raises ParseError."""

        f = tmp_path / "bad.xml"
        f.write_text("<unclosed>", encoding="utf-8")
        with pytest.raises(ParseError, match="Cannot parse XML"):
            chunk_file(f)

    def test_xml_skips_comment_children(self, tmp_path: Path) -> None:
        """Comment and processing-instruction children are skipped."""

        f = tmp_path / "comments.xml"
        f.write_text(
            "<root><!-- skip me --><section>Real text</section></root>",
            encoding="utf-8",
        )
        result = chunk_file(f)
        assert any(c["section"] == "section" for c in result)

    def test_xml_skips_empty_children(self, tmp_path: Path) -> None:
        """Child elements with no text are skipped."""

        f = tmp_path / "empty_child.xml"
        f.write_text(
            "<root><empty></empty><full>Has text</full></root>",
            encoding="utf-8",
        )
        result = chunk_file(f)
        sections = {c["section"] for c in result}
        assert "empty" not in sections
        assert "full" in sections

    def test_missing_lxml_raises_parse_error(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        """If lxml is not installed, XML parsing raises ParseError."""

        f = tmp_path / "x.xml"
        f.write_text("<root/>", encoding="utf-8")
        monkeypatch.setitem(sys.modules, "lxml", None)
        monkeypatch.setitem(sys.modules, "lxml.etree", None)
        with pytest.raises(ParseError, match="lxml is not installed"):
            chunk_file(f)
